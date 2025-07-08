import gradio as gr
from PIL import Image
import os, tempfile, uuid, traceback
from io import BytesIO
import docx, openpyxl
from pptx import Presentation

from google import genai
from google.genai import types

# --- Gemini API Configuration ---
GEMINI_API_KEY = "YOUR_API_KEY"

if not GEMINI_API_KEY:
    raise ValueError("Thiếu Gemini API Key. Mở thư mục app chạy: export GEMINI_API_KEY=...")

gemini = genai.Client(api_key=GEMINI_API_KEY)

# --- Global Defaults ---
DEFAULT_MAX_TOKEN = 500
DEFAULT_TEMPERATURE = 0.5
DEFAULT_TOPP = 1.0
DEFAULT_STOP = ""

TXT_MODEL = "gemini-1.5-flash"
VISION_MODEL = "gemini-2.0-flash"
IMG_MODEL = "gemini-2.0-flash-preview-image-generation"

def read_document_text(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".docx":
        doc = docx.Document(filepath)
        return "\n".join([p.text for p in doc.paragraphs])
    elif ext == ".xlsx":
        wb = openpyxl.load_workbook(filepath, data_only=True)
        t = ""
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                t += "\t".join([str(c) if c else "" for c in row]) + "\n"
        return t
    elif ext == ".pptx":
        prs = Presentation(filepath)
        t = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    t += shape.text + "\n"
        return t
    else:
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

def chat_stream(message, history, max_tokens, temperature, top_p, stop_seq, system_instruction):
    history = history or []
    history.append((message, None))
    yield "", history

    try:
        stops = [s.strip() for s in stop_seq.split(",") if s.strip()]
        stream = gemini.models.generate_content_stream(
            model=TXT_MODEL,
            config=types.GenerateContentConfig(
                system_instruction=system_instruction,
                max_output_tokens=max_tokens,
                temperature=temperature,
                top_p=top_p,
                stop_sequences=stops
            ),
            contents=[types.Content(role="user", parts=[types.Part(text=message)])]
        )
        reply = ""
        for chunk in stream:
            if chunk.text:
                reply += chunk.text
                history[-1] = (message, reply)
                yield "", history
    except Exception as e:
        traceback.print_exc()
        history[-1] = (message, f"Lỗi: {e}")
        yield "", history
        
def summarize_document(file_obj, history, *_):
    history = history or []
    try:
        ext = os.path.splitext(file_obj.name)[1].lower()
        if ext == ".pdf":
            with open(file_obj.name, "rb") as f:
                doc_data = f.read()
            resp = gemini.models.generate_content(
                model=VISION_MODEL,
                contents=[
                    types.Part.from_bytes(data=doc_data, mime_type="application/pdf"),
                    "Hãy tóm tắt tài liệu này"
                ]
            )
        else:
            text = read_document_text(file_obj.name)
            resp = gemini.models.generate_content(
                model=TXT_MODEL,
                contents=[
                    types.Part(text=text),
                    "Hãy tóm tắt tài liệu này"
                ]
            )
        history.append(("📄 Tóm tắt:", resp.text))
    except Exception as e:
        traceback.print_exc()
        history.append(("❌ Lỗi tóm tắt:", str(e)))
    return history

def describe_image(image, chatbox, max_tokens, temp, top_p, sys_inst):
    try:
        img_bytes = BytesIO()
        Image.fromarray(image).save(img_bytes, format="JPEG")
        response = gemini.models.generate_content(
            model=VISION_MODEL,
            config=types.GenerateContentConfig(
                system_instruction=sys_inst,
                temperature=temp,
                max_output_tokens=max_tokens,
                top_p=top_p,
            ),
            contents=[
                types.Part.from_bytes(data=img_bytes.getvalue(), mime_type="image/jpeg"),
                types.Part(text="Mô tả chi tiết hình ảnh trên bằng tiếng Việt."),
            ],
        )
        chatbox.append(("🖼️", response.text))
        return chatbox
    except Exception as e:
        traceback.print_exc()
        chatbox.append(("❌", f"Lỗi mô tả ảnh: {e}"))
        return chatbox
    
def generate_or_edit_image(prompt, image, history):
    history = history or []
    try:
        parts = []
        if image is not None:
            buf = BytesIO()
            Image.fromarray(image).save(buf, format="PNG")
            parts.append(types.Part.from_bytes(data=buf.getvalue(), mime_type="image/png"))
        parts.append(prompt)

        resp = gemini.models.generate_content(
            model=IMG_MODEL,
            config=types.GenerateContentConfig(response_modalities=["TEXT","IMAGE"]),
            contents=parts
        )
        img_path = None
        text = ""
        for p in resp.candidates[0].content.parts:
            if p.text:
                text = p.text
            if p.inline_data:
                img = Image.open(BytesIO(p.inline_data.data))
                img_path = os.path.join(tempfile.gettempdir(), f"gen_{uuid.uuid4().hex}.png")
                img.save(img_path)
        history.append(("🎨 Prompt:", prompt))
        if img_path: history.append(("🖼️ Ảnh:", {"path": img_path}))
        if text: history.append(("📝 Mô tả:", text))
    except Exception as e:
        traceback.print_exc()
        history.append(("❌ Lỗi tạo/chỉnh ảnh:", str(e)))
    return "", None, history

with gr.Blocks() as demo:
    gr.Markdown("# 🤖 Gemini Multimodal Chatbot")

    with gr.Accordion("⚙️ Cài đặt Chung", open=False):
        max_tok = gr.Slider(1, 65000, value=DEFAULT_MAX_TOKEN, label="Max Tokens")
        temp = gr.Slider(0,1, value=DEFAULT_TEMPERATURE, label="Temperature")
        top_p = gr.Slider(0,1, value=DEFAULT_TOPP, label="Top P")
        stop_s = gr.Textbox(value=DEFAULT_STOP, label="Stop Sequences (phân cách dấu ',')")

    with gr.Tabs():
        with gr.TabItem("💬 Chat"):
            sys_c = gr.Textbox(value="Bạn là trợ lý AI thân thiện, trả lời bằng tiếng Việt.", label="System instruction")
            chat_c = gr.Chatbot(layout="bubble")
            txt = gr.Textbox(placeholder="Nhập câu hỏi...", label=None)
            btn = gr.Button("Gửi")
            clr = gr.Button("🗑️ Xóa")

            btn.click(chat_stream, [txt, chat_c, max_tok, temp, top_p, stop_s, sys_c], [txt, chat_c])
            txt.submit(chat_stream, [txt, chat_c, max_tok, temp, top_p, stop_s, sys_c], [txt, chat_c])
            clr.click(lambda: [], None, chat_c)

        with gr.TabItem("📄 Tóm tắt Tài liệu"):
            sys_d = gr.Textbox(value="Bạn là trợ lý tóm tắt bằng tiếng Việt.", label="System instruction")
            chat_d = gr.Chatbot()
            up_d = gr.File(file_types=[".pdf",".docx",".xlsx",".txt"])
            clr_d = gr.Button("🗑️ Xóa")

            up_d.upload(summarize_document, [up_d, chat_d, max_tok, temp, top_p, sys_d],[chat_d])
            clr_d.click(lambda: [], None, chat_d)

        with gr.TabItem("🖼️ Mô tả Ảnh"):
            sys_i = gr.Textbox(value="Bạn là trợ lý mô tả ảnh, tiếng Việt.", label="System instruction")
            chat_i = gr.Chatbot()
            up_i = gr.Image(type="numpy")
            clr_i = gr.Button("🗑️ Xóa")

            up_i.upload(describe_image, [up_i, chat_i, max_tok, temp, top_p, sys_i],[chat_i])
            clr_i.click(lambda: [], None, chat_i)

        with gr.TabItem("🎨 Tạo/Chỉnh Ảnh"):
            sys_g = gr.Textbox(value="Tạo hoặc chỉnh sửa ảnh với prompt tiếng Việt.", label="Prompt")
            img_g = gr.Image(type="numpy")
            chat_g = gr.Chatbot()
            btn_g = gr.Button("Tạo/Chỉnh")
            clr_g = gr.Button("🗑️ Xóa")

            btn_g.click(generate_or_edit_image, [sys_g, img_g, chat_g],[sys_g, img_g, chat_g])
            clr_g.click(lambda: [], None, chat_g)

    demo.launch()
